from io import BytesIO
from pptx import Presentation
from fastapi import FastAPI, File, Response, UploadFile
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageOps

app = FastAPI()


@app.post("/upload")
def upload(
    upload: UploadFile = File(
        description="PowerPoint file to be uploaded and inverted",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ),
):
    try:
        contents = upload.file.read()
        bytes = BytesIO(contents)

        presentation = Presentation(bytes)

        for slide in presentation.slides:
            # Invert slide background to black
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # Invert text color to white for all shapes that contain text
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract the image from the shape
                    image_stream = BytesIO(shape.image.blob)
                    with Image.open(image_stream) as img:
                        # If image format is not supported convert it first
                        if img.mode not in ("RGB", "RGBA"):
                            img = img.convert("RGB")

                        # Invert the image color using Pillow only if supported
                        inverted_img = ImageOps.invert(img)

                        # Save the inverted image to a stream
                        img_stream = BytesIO()
                        inverted_img.save(img_stream, format="PNG")
                        img_stream.seek(0)

                        # Remove the original picture
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        slide.shapes._spTree.remove(shape._element)

                        # Add the new inverted image to the slide
                        slide.shapes.add_picture(img_stream, left, top, width, height)

                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)

        output = BytesIO()
        presentation.save(output)

        return Response(
            content=output.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        # with open(file.filename, "wb") as f:
        #     f.write(contents)
    except Exception as e:
        print(e)
        return {"message": "There was an error uploading the file"}
    finally:
        upload.file.close()
